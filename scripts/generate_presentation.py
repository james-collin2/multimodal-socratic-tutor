"""
scripts/generate_presentation.py

Generates docs/socratOT_presentation.pptx — a dark-themed deck matching the app.
Run: python scripts/generate_presentation.py
"""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parent.parent
OUT = ROOT / "docs" / "socratOT_presentation.pptx"

# ── Palette (matches the app's Dark Professional theme) ──────────────────────
BG = RGBColor(0x0B, 0x11, 0x20)
SURFACE = RGBColor(0x15, 0x1D, 0x2E)
SURFACE2 = RGBColor(0x1A, 0x23, 0x36)
TEXT = RGBColor(0xE2, 0xE8, 0xF0)
MUTED = RGBColor(0x94, 0xA3, 0xB8)
INDIGO = RGBColor(0x81, 0x8C, 0xF8)
INDIGO_D = RGBColor(0x63, 0x66, 0xF1)
CYAN = RGBColor(0x22, 0xD3, 0xEE)
GREEN = RGBColor(0x6E, 0xE7, 0xB7)
AMBER = RGBColor(0xFB, 0xBF, 0x24)
BORDER = RGBColor(0x2A, 0x36, 0x50)

W, H = Inches(13.333), Inches(7.5)
FONT = "Inter"


def _bg(slide) -> None:
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = BG


def _box(slide, x, y, w, h, fill=SURFACE, line=BORDER, radius=True):
    from pptx.enum.shapes import MSO_SHAPE

    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE, x, y, w, h
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = line
    shape.line.width = Pt(1)
    shape.shadow.inherit = False
    return shape


def _text(slide, x, y, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, space=4):
    """runs: list of (text, size, color, bold) — each becomes a paragraph."""
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    for i, (txt, size, color, bold) in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_after = Pt(space)
        r = p.add_run()
        r.text = txt
        r.font.size = Pt(size)
        r.font.color.rgb = color
        r.font.bold = bold
        r.font.name = FONT
    return tb


def _accent_bar(slide, x, y, w=None):
    bar = _box(slide, x, y, w or Inches(0.55), Inches(0.07), fill=CYAN, line=CYAN)
    return bar


def _header(slide, kicker, title):
    _accent_bar(slide, Inches(0.7), Inches(0.62))
    _text(slide, Inches(0.7), Inches(0.72), Inches(12), Inches(0.4), [(kicker, 13, CYAN, True)])
    _text(slide, Inches(0.7), Inches(1.0), Inches(12), Inches(0.9), [(title, 30, TEXT, True)])


def _slide():
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(s)
    return s


prs = Presentation()
prs.slide_width = W
prs.slide_height = H

# ── 1. Title ─────────────────────────────────────────────────────────────────
s = _slide()
_box(s, Inches(0.6), Inches(2.35), Inches(0.9), Inches(0.9), fill=INDIGO_D, line=INDIGO_D)
_text(
    s,
    Inches(0.6),
    Inches(2.35),
    Inches(0.9),
    Inches(0.9),
    [("S", 40, RGBColor(255, 255, 255), True)],
    align=PP_ALIGN.CENTER,
    anchor=MSO_ANCHOR.MIDDLE,
)
_text(s, Inches(1.7), Inches(2.3), Inches(11), Inches(1.2), [("socratOT", 54, TEXT, True)])
_text(
    s,
    Inches(1.72),
    Inches(3.35),
    Inches(11),
    Inches(0.6),
    [("Socratic Multimodal RAG Tutor for Occupational Therapy Education", 20, INDIGO, False)],
)
_text(
    s,
    Inches(1.72),
    Inches(4.05),
    Inches(11),
    Inches(0.8),
    [
        (
            "Teaching Anatomy & Neuroscience through guided questioning — not direct answers.",
            15,
            MUTED,
            False,
        )
    ],
)
_accent_bar(s, Inches(1.75), Inches(4.75), w=Inches(2.0))
_text(
    s,
    Inches(1.72),
    Inches(4.95),
    Inches(11),
    Inches(0.5),
    [("Bahodir Nematjonov · Inha University · June 2026", 13, MUTED, False)],
)

# ── 2. The problem ───────────────────────────────────────────────────────────
s = _slide()
_header(s, "MOTIVATION", "Answering is not teaching")
_text(
    s,
    Inches(0.7),
    Inches(2.1),
    Inches(11.9),
    Inches(1.0),
    [
        (
            "LLMs answer anatomy questions fluently — but OT education depends on clinical "
            "reasoning: connecting structure → function → therapeutic implication.",
            18,
            TEXT,
            False,
        )
    ],
)
cards = [
    (
        "The gap",
        "A bot that states “the cerebellum coordinates movement” removes the reasoning practice that builds durable knowledge.",
        AMBER,
    ),
    (
        "The constraint",
        "The tutor must guide via questions and must NOT reveal the answer in the first two turns.",
        INDIGO,
    ),
    (
        "The goal",
        "A complete, grounded, multimodal platform built around that single pedagogical rule.",
        GREEN,
    ),
]
x = Inches(0.7)
for title, body, col in cards:
    _box(s, x, Inches(3.4), Inches(3.85), Inches(2.6))
    _accent_bar(s, x + Inches(0.3), Inches(3.7), w=Inches(0.5))
    _text(s, x + Inches(0.3), Inches(3.85), Inches(3.25), Inches(0.5), [(title, 18, col, True)])
    _text(s, x + Inches(0.3), Inches(4.4), Inches(3.25), Inches(1.5), [(body, 14, MUTED, False)])
    x += Inches(4.05)

# ── 3. System architecture ───────────────────────────────────────────────────
s = _slide()
_header(s, "ARCHITECTURE", "Nine integrated components")
comps = [
    "1 · Knowledge base & dataset",
    "2 · RAG pipeline (+ citations, hallucination guard)",
    "3 · Socratic logic (Tutor-Not-Teller)",
    "4 · Conversation state (4 phases)",
    "5 · Multimodal diagram tutoring",
    "6 · Cross-session student memory",
    "7 · Reasoning assessment engine",
    "8 · Student dashboard",
    "9 · Accessibility (TTS / STT)",
    "Eval · RAGAS + compliance + baselines",
]
x, y = Inches(0.7), Inches(2.1)
for i, c in enumerate(comps):
    col = i % 2
    row = i // 2
    bx = Inches(0.7) + col * Inches(6.1)
    by = Inches(2.15) + row * Inches(1.0)
    _box(s, bx, by, Inches(5.8), Inches(0.82), fill=SURFACE)
    _text(
        s,
        bx + Inches(0.3),
        by,
        Inches(5.3),
        Inches(0.82),
        [(c, 15, TEXT, False)],
        anchor=MSO_ANCHOR.MIDDLE,
    )

# ── 4. Knowledge base & dataset ──────────────────────────────────────────────
s = _slide()
_header(s, "REQUIREMENT 1", "Knowledge base & dataset")
stats = [
    ("2,322", "text chunks (OpenStax A&P 2e)"),
    ("723", "anatomical images + JSON metadata"),
    ("50", "ground-truth Q&A pairs"),
    ("1,536-d", "OpenAI embeddings · ChromaDB"),
]
x = Inches(0.7)
for big, small in stats:
    _box(s, x, Inches(2.3), Inches(2.85), Inches(1.7))
    _text(
        s,
        x,
        Inches(2.55),
        Inches(2.85),
        Inches(0.7),
        [(big, 34, INDIGO, True)],
        align=PP_ALIGN.CENTER,
    )
    _text(
        s,
        x + Inches(0.2),
        Inches(3.35),
        Inches(2.45),
        Inches(0.6),
        [(small, 12, MUTED, False)],
        align=PP_ALIGN.CENTER,
    )
    x += Inches(3.0)
_text(
    s,
    Inches(0.7),
    Inches(4.5),
    Inches(11.9),
    Inches(2),
    [
        (
            "Corpus: OpenStax Anatomy & Physiology 2e (CC BY 4.0), cleaned and chunked (512 / 64 overlap).",
            15,
            TEXT,
            False,
        ),
        (
            "Images: 723 figures (≈459 MB) catalogued with page, dimensions, source, license — exceeds the 20–30 requirement.",
            15,
            TEXT,
            False,
        ),
        (
            "Ground truth: 45 anatomy + 5 deliberately out-of-corpus clinical questions (to test faithful abstention).",
            15,
            TEXT,
            False,
        ),
    ],
)

# ── 5. RAG pipeline ──────────────────────────────────────────────────────────
s = _slide()
_header(s, "REQUIREMENT 2", "RAG with grounding & hallucination control")
flow = [
    "Query",
    "Retriever",
    "ChromaDB\n(top-k=8)",
    "LLM\n(gpt-4o-mini)",
    "Hallucination\nguard",
    "Grounded\nanswer + cites",
]
x = Inches(0.7)
for i, step in enumerate(flow):
    _box(s, x, Inches(2.5), Inches(1.7), Inches(1.2), fill=SURFACE2)
    _text(
        s,
        x,
        Inches(2.5),
        Inches(1.7),
        Inches(1.2),
        [(step, 13, TEXT, True)],
        align=PP_ALIGN.CENTER,
        anchor=MSO_ANCHOR.MIDDLE,
    )
    if i < len(flow) - 1:
        _text(
            s,
            x + Inches(1.7),
            Inches(2.5),
            Inches(0.32),
            Inches(1.2),
            [("→", 20, CYAN, True)],
            align=PP_ALIGN.CENTER,
            anchor=MSO_ANCHOR.MIDDLE,
        )
    x += Inches(2.02)
_text(
    s,
    Inches(0.7),
    Inches(4.3),
    Inches(11.9),
    Inches(2),
    [
        (
            "Two-stage hallucination guard: fast keyword-overlap check → LLM semantic verification.",
            16,
            TEXT,
            False,
        ),
        (
            "Every answer carries source citations (e.g. openStax/anatomy_physiology_2e.pdf#pp_313–320).",
            16,
            TEXT,
            False,
        ),
        (
            "Out-of-corpus questions → the system faithfully abstains instead of fabricating.",
            16,
            GREEN,
            False,
        ),
    ],
)

# ── 6. Socratic engine ───────────────────────────────────────────────────────
s = _slide()
_header(s, "REQUIREMENT 3 — CRITICAL", "The Tutor-Not-Teller engine")
_text(
    s,
    Inches(0.7),
    Inches(2.05),
    Inches(11.9),
    Inches(0.5),
    [("Four-phase state machine: rapport → tutoring → assessment → mastery", 16, INDIGO, True)],
)
phases = [
    ("NONE", "Guiding question"),
    ("HINT 1", "Open nudge"),
    ("HINT 2", "Pointed clue"),
    ("REVEAL", "Answer — only now"),
]
x = Inches(0.7)
for i, (lvl, desc) in enumerate(phases):
    col = GREEN if i == 3 else INDIGO
    _box(s, x, Inches(2.7), Inches(2.7), Inches(1.3))
    _text(
        s, x, Inches(2.85), Inches(2.7), Inches(0.5), [(lvl, 18, col, True)], align=PP_ALIGN.CENTER
    )
    _text(
        s,
        x,
        Inches(3.4),
        Inches(2.7),
        Inches(0.5),
        [(desc, 13, MUTED, False)],
        align=PP_ALIGN.CENTER,
    )
    if i < 3:
        _text(
            s,
            x + Inches(2.7),
            Inches(2.7),
            Inches(0.3),
            Inches(1.3),
            [("→", 18, CYAN, True)],
            align=PP_ALIGN.CENTER,
            anchor=MSO_ANCHOR.MIDDLE,
        )
    x += Inches(2.95)
_text(
    s,
    Inches(0.7),
    Inches(4.5),
    Inches(11.9),
    Inches(2),
    [
        (
            "•  Reveal is gated by hint-level ≥ 2 — structurally no direct answer in the first two turns.",
            15,
            TEXT,
            False,
        ),
        (
            "•  Knowledge masking: the LLM references retrieved context without stating the target fact.",
            15,
            TEXT,
            False,
        ),
        (
            "•  Bypass resistance: 23-pattern detector blocks “just tell me” / prompt-injection attempts.",
            15,
            TEXT,
            False,
        ),
    ],
)

# ── 7. Multimodal ────────────────────────────────────────────────────────────
s = _slide()
_header(s, "REQUIREMENT 5", "Multimodal diagram tutoring")
_text(
    s,
    Inches(0.7),
    Inches(2.1),
    Inches(11.9),
    Inches(1.0),
    [
        (
            "Upload any anatomy diagram — including unseen images (blind testing). GPT-4o vision identifies "
            "structures; they are routed back into the RAG corpus to keep follow-up questioning grounded.",
            17,
            TEXT,
            False,
        )
    ],
)
steps = [
    "Image upload",
    "GPT-4o vision\n(structures, region, confidence)",
    "RAG grounding\non structures",
    "3 Socratic questions\n(graded by difficulty)",
]
x = Inches(0.7)
for i, st_ in enumerate(steps):
    _box(s, x, Inches(3.5), Inches(2.7), Inches(1.4), fill=SURFACE2)
    _text(
        s,
        x,
        Inches(3.5),
        Inches(2.7),
        Inches(1.4),
        [(st_, 13, TEXT, True)],
        align=PP_ALIGN.CENTER,
        anchor=MSO_ANCHOR.MIDDLE,
    )
    if i < 3:
        _text(
            s,
            x + Inches(2.7),
            Inches(3.5),
            Inches(0.3),
            Inches(1.4),
            [("→", 18, CYAN, True)],
            align=PP_ALIGN.CENTER,
            anchor=MSO_ANCHOR.MIDDLE,
        )
    x += Inches(2.95)
_text(
    s,
    Inches(0.7),
    Inches(5.3),
    Inches(11.9),
    Inches(1),
    [
        (
            "Live example: a reproductive-system diagram → 12 structures, region=Uterus, 95% confidence, "
            "3 questions + OpenStax citations.",
            14,
            GREEN,
            False,
        )
    ],
)

# ── 8. Memory & assessment ───────────────────────────────────────────────────
s = _slide()
_header(s, "REQUIREMENTS 6 – 8", "Memory · assessment · dashboard")
items = [
    (
        "Student memory",
        "SQLite cross-session store of weak/strong topics & mastery scores. Revisits missed concepts: “Last session you struggled with Cranial Nerve VII — let's revisit it.”",
        INDIGO,
    ),
    (
        "Reasoning assessment",
        "LLM-as-judge scores clinical reasoning 0–100 across clinical accuracy, reasoning quality, and terminology — not multiple choice.",
        CYAN,
    ),
    (
        "Student dashboard",
        "Topic mastery, weak areas, performance trends, session history & topic table — fed live by the tutoring and image flows.",
        GREEN,
    ),
]
y = Inches(2.2)
for title, body, col in items:
    _box(s, Inches(0.7), y, Inches(11.9), Inches(1.45))
    _accent_bar(s, Inches(1.0), y + Inches(0.28), w=Inches(0.5))
    _text(s, Inches(1.0), y + Inches(0.42), Inches(3.2), Inches(0.8), [(title, 18, col, True)])
    _text(
        s,
        Inches(4.3),
        y + Inches(0.18),
        Inches(8.0),
        Inches(1.1),
        [(body, 14, TEXT, False)],
        anchor=MSO_ANCHOR.MIDDLE,
    )
    y += Inches(1.6)

# ── 9. Evaluation: compliance ────────────────────────────────────────────────
s = _slide()
_header(s, "EVALUATION 1", "Socratic compliance — perfect")
res = [
    ("100%", "Bypass detection (20/20)", GREEN),
    ("0%", "False positives (0/12)", GREEN),
    ("100%", "Socratic compliance", GREEN),
]
x = Inches(1.3)
for big, small, col in res:
    _box(s, x, Inches(2.6), Inches(3.3), Inches(2.1))
    _text(
        s, x, Inches(2.95), Inches(3.3), Inches(0.9), [(big, 48, col, True)], align=PP_ALIGN.CENTER
    )
    _text(
        s,
        x + Inches(0.2),
        Inches(3.95),
        Inches(2.9),
        Inches(0.6),
        [(small, 14, MUTED, False)],
        align=PP_ALIGN.CENTER,
    )
    x += Inches(3.5)
_text(
    s,
    Inches(0.7),
    Inches(5.1),
    Inches(11.9),
    Inches(1),
    [
        (
            "Every adversarial “just tell me the answer” / instruction-injection attempt was caught and "
            "redirected to a hint — with zero legitimate questions wrongly blocked.",
            15,
            TEXT,
            False,
        )
    ],
)

# ── 10. Evaluation: RAGAS ────────────────────────────────────────────────────
s = _slide()
_header(s, "EVALUATION 2", "RAGAS — faithfulness 0.88")
bars = [
    ("Faithfulness", 0.88, GREEN),
    ("Answer relevance", 0.722, INDIGO),
    ("Context recall", 0.583, AMBER),
]
y = Inches(2.6)
for label, val, col in bars:
    _text(
        s,
        Inches(0.7),
        y,
        Inches(3.2),
        Inches(0.5),
        [(label, 16, TEXT, True)],
        anchor=MSO_ANCHOR.MIDDLE,
    )
    _box(s, Inches(4.0), y + Inches(0.05), Inches(7.5), Inches(0.45), fill=SURFACE2, line=BORDER)
    _box(s, Inches(4.0), y + Inches(0.05), Inches(7.5 * val), Inches(0.45), fill=col, line=col)
    _text(
        s,
        Inches(11.7),
        y,
        Inches(1.2),
        Inches(0.5),
        [(f"{val:.3f}", 16, col, True)],
        anchor=MSO_ANCHOR.MIDDLE,
    )
    y += Inches(0.85)
_text(
    s,
    Inches(0.7),
    Inches(5.4),
    Inches(11.9),
    Inches(1),
    [
        (
            "High faithfulness: answers stay grounded, and the system abstains on out-of-corpus questions. "
            "Context recall is the clearest target for future retrieval tuning (reranking, query expansion).",
            15,
            TEXT,
            False,
        )
    ],
)

# ── 11. Evaluation: baselines / key insight ──────────────────────────────────
s = _slide()
_header(s, "EVALUATION 3", "Baselines & the key insight")
rows = [
    ("System", "ROUGE-L", "BERTScore", "Overlap", True),
    ("socratOT (full)", "0.225", "0.257", "0.232", False),
    ("no-RAG", "0.351", "0.416", "0.473", False),
    ("no-Socratic", "0.278", "0.338", "0.526", False),
]
y = Inches(2.15)
for name, a, b, c, head in rows:
    fill = SURFACE2 if head else SURFACE
    _box(s, Inches(0.7), y, Inches(7.0), Inches(0.62), fill=fill)
    cols = [
        (name, Inches(0.9), 2.7),
        (a, Inches(4.0), 1.0),
        (b, Inches(5.3), 1.2),
        (c, Inches(6.7), 1.0),
    ]
    for txt, cx, cw in cols:
        _text(
            s,
            cx,
            y,
            Inches(cw),
            Inches(0.62),
            [(txt, 14, TEXT if head else MUTED, head)],
            anchor=MSO_ANCHOR.MIDDLE,
        )
    y += Inches(0.66)
_box(s, Inches(8.1), Inches(2.15), Inches(4.5), Inches(3.2), fill=SURFACE, line=INDIGO_D)
_text(s, Inches(8.4), Inches(2.4), Inches(3.9), Inches(0.5), [("The key finding", 18, CYAN, True)])
_text(
    s,
    Inches(8.4),
    Inches(2.95),
    Inches(3.9),
    Inches(2.4),
    [
        ("socratOT scores LOWEST on overlap metrics — on purpose.", 15, TEXT, True),
        (
            "Overlap rewards restating the reference answer; the Socratic tutor deliberately withholds it.",
            13,
            MUTED,
            False,
        ),
        (
            "→ Compliance & faithfulness are the meaningful axes for tutoring systems.",
            13,
            GREEN,
            True,
        ),
    ],
)

# ── 12. Conclusion ───────────────────────────────────────────────────────────
s = _slide()
_header(s, "CONCLUSION", "A behaviour-aware tutoring platform")
points = [
    ("✓", "All 9 requirements integrated end-to-end — not a chatbot.", GREEN),
    ("✓", "Perfect Socratic compliance + 0.88 faithfulness, with grounded citations.", GREEN),
    (
        "✓",
        "Multimodal diagram tutoring on unseen images, cross-session memory, reasoning assessment.",
        GREEN,
    ),
    (
        "→",
        "Future work: local open-source inference, retrieval reranking, longitudinal learning study.",
        CYAN,
    ),
]
y = Inches(2.4)
for mark, txt, col in points:
    _text(
        s,
        Inches(0.9),
        y,
        Inches(0.6),
        Inches(0.6),
        [(mark, 22, col, True)],
        anchor=MSO_ANCHOR.MIDDLE,
    )
    _text(
        s,
        Inches(1.6),
        y,
        Inches(11),
        Inches(0.7),
        [(txt, 18, TEXT, False)],
        anchor=MSO_ANCHOR.MIDDLE,
    )
    y += Inches(0.95)
_accent_bar(s, Inches(0.9), Inches(6.3), w=Inches(2.0))
_text(
    s,
    Inches(0.9),
    Inches(6.45),
    Inches(11),
    Inches(0.5),
    [("socratOT — teaching reasoning, not answers.", 16, INDIGO, True)],
)


def main() -> None:
    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUT))
    print(f"Saved {len(prs.slides._sldIdLst)} slides → {OUT}")


if __name__ == "__main__":
    main()
